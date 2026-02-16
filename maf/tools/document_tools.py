"""Document generation tools — create XLSX and PPTX from templates.

Uses openpyxl for Excel and python-pptx for PowerPoint.  Generated files
are written to a local `output/` directory and can later be uploaded to
SharePoint via the SharePoint agent.
"""

from __future__ import annotations

import json
import os
from datetime import datetime
from pathlib import Path
from typing import Annotated

from agent_framework import tool

# Resolve paths relative to this file
_BASE_DIR = Path(__file__).resolve().parent.parent
_TEMPLATES_DIR = _BASE_DIR / "templates"
_OUTPUT_DIR = _BASE_DIR / "output"


def _ensure_output_dir() -> Path:
    _OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return _OUTPUT_DIR


# ── XLSX Generation ──────────────────────────────────────────────────────

@tool(approval_mode="never_require")
def create_xlsx_report(
    title: Annotated[str, "Report title that appears in the header"],
    data_json: Annotated[
        str,
        "JSON string with 'headers' (list of column names) and 'rows' (list of lists). "
        "Example: {\"headers\":[\"Task\",\"Owner\",\"Status\"],\"rows\":[[\"Auth\",\"Bob\",\"Done\"]]}"
    ],
) -> str:
    """Create an Excel report (.xlsx) from the provided data and return the file path.

    The report uses the built-in project_status template with professional formatting.
    """
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    except ImportError as exc:
        return json.dumps({"error": "openpyxl is not installed. Run: pip install openpyxl"})

    data = json.loads(data_json)
    headers = data.get("headers", [])
    rows = data.get("rows", [])

    wb = Workbook()
    ws = wb.active
    ws.title = "Project Report"

    # ── Title row ────────────────────────────────────────────────────
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(len(headers), 3))
    title_cell = ws.cell(row=1, column=1, value=title)
    title_cell.font = Font(name="Calibri", size=16, bold=True, color="1F4E79")
    title_cell.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 40

    # ── Date row ─────────────────────────────────────────────────────
    ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=max(len(headers), 3))
    date_cell = ws.cell(row=2, column=1, value=f"Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}")
    date_cell.font = Font(name="Calibri", size=10, italic=True, color="666666")
    date_cell.alignment = Alignment(horizontal="center")

    # ── Header row ───────────────────────────────────────────────────
    header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    thin_border = Border(
        left=Side(style="thin"),
        right=Side(style="thin"),
        top=Side(style="thin"),
        bottom=Side(style="thin"),
    )

    for col_idx, header in enumerate(headers, start=1):
        cell = ws.cell(row=4, column=col_idx, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = thin_border

    # ── Data rows ────────────────────────────────────────────────────
    alt_fill = PatternFill(start_color="D6E4F0", end_color="D6E4F0", fill_type="solid")
    data_font = Font(name="Calibri", size=11)

    for row_idx, row_data in enumerate(rows, start=5):
        for col_idx, value in enumerate(row_data, start=1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            cell.font = data_font
            cell.border = thin_border
            cell.alignment = Alignment(vertical="center")
            if (row_idx - 5) % 2 == 1:
                cell.fill = alt_fill

    # ── Auto-fit column widths ───────────────────────────────────────
    for col_idx in range(1, len(headers) + 1):
        max_len = max(
            len(str(ws.cell(row=r, column=col_idx).value or ""))
            for r in range(4, 5 + len(rows))
        )
        ws.column_dimensions[ws.cell(row=4, column=col_idx).column_letter].width = max(max_len + 4, 12)

    # ── Save ─────────────────────────────────────────────────────────
    output_dir = _ensure_output_dir()
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"report_{ts}.xlsx"
    filepath = output_dir / filename
    wb.save(str(filepath))

    return json.dumps({
        "status": "success",
        "file_path": str(filepath.relative_to(_BASE_DIR.parent)),
        "file_name": filename,
        "sheets": ["Project Report"],
        "row_count": len(rows),
    }, indent=2)


# ── PPTX Generation ─────────────────────────────────────────────────────

@tool(approval_mode="never_require")
def create_pptx_presentation(
    title: Annotated[str, "Presentation title for the title slide"],
    subtitle: Annotated[str, "Subtitle or project name"] = "Project Summary",
    slides_json: Annotated[
        str,
        "JSON array of slide objects. Each object has 'title' and 'bullets' (list of strings). "
        "Example: [{\"title\":\"Status\",\"bullets\":[\"On track\",\"2 risks\"]}]"
    ] = "[]",
) -> str:
    """Create a PowerPoint presentation (.pptx) with a title slide and content slides.

    Returns the path to the generated file.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return json.dumps({"error": "python-pptx is not installed. Run: pip install python-pptx"})

    slides_data = json.loads(slides_json)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Colour palette ───────────────────────────────────────────────
    DARK_BG = RGBColor(0x1F, 0x4E, 0x79)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    ACCENT = RGBColor(0x2E, 0x86, 0xC1)
    LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)

    # ── Title Slide ──────────────────────────────────────────────────
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title text
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT
    p2.alignment = PP_ALIGN.CENTER

    # Date line
    p3 = tf.add_paragraph()
    p3.text = datetime.now().strftime("%B %d, %Y")
    p3.font.size = Pt(14)
    p3.font.color.rgb = WHITE
    p3.font.italic = True
    p3.alignment = PP_ALIGN.CENTER

    # ── Content Slides ───────────────────────────────────────────────
    for slide_info in slides_data:
        s = prs.slides.add_slide(prs.slide_layouts[6])

        # Accent bar at top
        from pptx.util import Inches as In
        shape = s.shapes.add_shape(
            1,  # Rectangle
            In(0), In(0), prs.slide_width, In(0.6),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_BG
        shape.line.fill.background()

        # Slide title in the accent bar
        title_box = s.shapes.add_textbox(In(0.5), In(0.05), In(12), In(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_info.get("title", "")
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Bullet points
        bullets = slide_info.get("bullets", [])
        if bullets:
            body_box = s.shapes.add_textbox(In(1), In(1.2), In(11), In(5.5))
            tf = body_box.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"•  {bullet}"
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                p.space_after = Pt(12)

    # ── Save ─────────────────────────────────────────────────────────
    output_dir = _ensure_output_dir()
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"presentation_{ts}.pptx"
    filepath = output_dir / filename
    prs.save(str(filepath))

    return json.dumps({
        "status": "success",
        "file_path": str(filepath.relative_to(_BASE_DIR.parent)),
        "file_name": filename,
        "slide_count": 1 + len(slides_data),
    }, indent=2)
